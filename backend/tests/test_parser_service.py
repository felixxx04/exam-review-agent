"""Tests for the ParserService with PDF, Word, PPT, and image support."""

import pytest

from app.services.parser_service import Chunk, ParserService


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

SAMPLE_TEXT = """\
第一章 初等几何

第一节 基本概念

几何学是研究空间形状和大小关系的学科。初等几何主要研究平面
图形和立体图形的基本性质。

定义1：点是最基本的几何元素，没有大小，只有位置。
定义2：线是由无数个点组成的集合，具有长度但没有宽度。

第二节 角度与三角形

角度是由两条从同一点出发的射线组成的图形。三角形是由三条
线段首尾相连组成的封闭图形。

定理1：三角形内角和等于一百八十度。

证明：设三角形ABC，过点A作直线DE平行于BC。
则∠DAB = ∠ABC，∠EAC = ∠ACB（内错角相等）。
因为∠DAB + ∠BAC + ∠EAC = 180°，所以∠ABC + ∠BAC + ∠ACB = 180°。

第三节 四边形

四边形是由四条线段首尾相连组成的封闭图形。常见的四边形包括
平行四边形、矩形、菱形、正方形和梯形。
"""


@pytest.fixture
def sample_pdf_path(tmp_path):
    """Create a real sample PDF using PyMuPDF for testing.

    Uses English text because PyMuPDF's default built-in fonts
    do not include CJK glyphs. Chinese text is tested via DOCX/PPTX.
    """
    import fitz

    path = tmp_path / "sample.pdf"
    doc = fitz.open()
    # First page
    page = doc.new_page()
    page.insert_text(
        (50, 50),
        "Chapter 1: Elementary Geometry\n\nSection 1.1: Basic Concepts\n\n"
        "Geometry is the study of spatial shapes and their size relationships.\n"
        "Elementary geometry primarily studies the basic properties of\n"
        "planar figures and solid figures.\n\n"
        "Definition 1: A point is the most fundamental geometric element,\n"
        "having no size, only position.\n"
        "Definition 2: A line is a collection of infinitely many points,\n"
        "having length but no width.",
        fontsize=12,
    )
    # Second page
    page2 = doc.new_page()
    page2.insert_text(
        (50, 50),
        "Section 1.2: Angles and Triangles\n\n"
        "An angle is a figure formed by two rays from the same point.\n"
        "A triangle is a closed figure formed by three line segments\n"
        "connected end to end.\n\n"
        "Theorem 1: The sum of the interior angles of a triangle equals 180 degrees.\n\n"
        "Proof: Given triangle ABC, draw line DE through point A parallel to BC.\n"
        "Then angle DAB = angle ABC, angle EAC = angle ACB\n"
        "(alternate interior angles are equal).\n"
        "Since DAB + BAC + EAC = 180 degrees,\n"
        "we have ABC + BAC + ACB = 180 degrees.",
        fontsize=12,
    )
    # Third page
    page3 = doc.new_page()
    page3.insert_text(
        (50, 50),
        "Section 1.3: Quadrilaterals\n\n"
        "A quadrilateral is a closed figure formed by four line segments\n"
        "connected end to end. Common quadrilaterals include\n"
        "parallelograms, rectangles, rhombuses, squares, and trapezoids.",
        fontsize=12,
    )
    doc.save(str(path))
    doc.close()
    return str(path)


@pytest.fixture
def sample_docx_path(tmp_path):
    """Create a sample .docx file for testing."""
    from docx import Document

    path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_heading("第一章 初等几何", level=1)
    doc.add_heading("第一节 基本概念", level=2)
    doc.add_paragraph(
        "几何学是研究空间形状和大小关系的学科。"
        "初等几何主要研究平面图形和立体图形的基本性质。"
    )
    doc.add_paragraph(
        "定义1：点是最基本的几何元素，没有大小，只有位置。"
    )
    doc.add_paragraph(
        "定义2：线是由无数个点组成的集合，具有长度但没有宽度。"
    )
    doc.add_heading("第二节 角度与三角形", level=2)
    doc.add_paragraph(
        "角度是由两条从同一点出发的射线组成的图形。"
    )
    doc.add_paragraph(
        "三角形是由三条线段首尾相连组成的封闭图形。"
    )
    doc.save(str(path))
    return str(path)


@pytest.fixture
def sample_pptx_path(tmp_path):
    """Create a sample .pptx file for testing."""
    from pptx import Presentation

    path = tmp_path / "sample.pptx"
    prs = Presentation()
    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "初等几何 - 基本概念"
    content1 = slide1.placeholders[1]
    content1.text = (
        "几何学是研究空间形状和大小关系的学科。\n"
        "点：最基本的几何元素，没有大小，只有位置。\n"
        "线：由无数个点组成的集合。"
    )
    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "三角形性质"
    content2 = slide2.placeholders[1]
    content2.text = (
        "三角形内角和等于一百八十度。\n"
        "证明：过点A作直线DE平行于BC。"
    )
    prs.save(str(path))
    return str(path)


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


class TestParserService:
    """Tests for the ParserService covering PDF, Word, PPT, and error handling."""

    async def test_parse_pdf_extracts_text_and_metadata(self, sample_pdf_path):
        """PDF parsing should extract text chunks with page metadata."""
        service = ParserService()
        result = await service.parse(sample_pdf_path, file_type="pdf")

        assert len(result.chunks) > 0
        assert result.page_count == 3
        # Check chunk has metadata
        first_chunk = result.chunks[0]
        assert "source" in first_chunk.metadata
        assert "sample" in first_chunk.metadata["source"]
        assert "page" in first_chunk.metadata
        # Check content
        text_content = " ".join(c.text for c in result.chunks)
        assert "Geometry" in text_content

    async def test_parse_pdf_chunks_have_increasing_indices(self, sample_pdf_path):
        """Chunk indices should be sequential and unique."""
        service = ParserService()
        result = await service.parse(sample_pdf_path, file_type="pdf")

        indices = [c.chunk_index for c in result.chunks]
        assert indices == list(range(len(indices)))

    async def test_parse_docx_extracts_paragraphs(self, sample_docx_path):
        """Word document parsing should extract text content."""
        service = ParserService()
        result = await service.parse(sample_docx_path, file_type="docx")

        assert len(result.chunks) > 0
        assert result.page_count > 0
        text_content = " ".join(c.text for c in result.chunks)
        assert "几何学" in text_content
        assert "定义" in text_content

    async def test_parse_pptx_extracts_slides(self, sample_pptx_path):
        """PowerPoint parsing should extract slide content."""
        service = ParserService()
        result = await service.parse(sample_pptx_path, file_type="pptx")

        assert len(result.chunks) > 0
        assert result.page_count == 2
        text_content = " ".join(c.text for c in result.chunks)
        assert "初等几何" in text_content
        assert "三角形" in text_content

    async def test_parse_unknown_type_raises_error(self, tmp_path):
        """Parsing an unsupported file type should raise ValueError."""
        path = tmp_path / "test.xyz"
        path.write_text("dummy")

        service = ParserService()
        with pytest.raises(ValueError, match="Unsupported file type"):
            await service.parse(str(path), file_type="xyz")

    async def test_parse_detects_file_type_from_extension(self, sample_pdf_path):
        """When file_type is not provided, detect from extension."""
        service = ParserService()
        result = await service.parse(sample_pdf_path)

        assert len(result.chunks) > 0
        assert result.page_count == 3

    async def test_parse_result_metadata_includes_file_info(self, sample_pdf_path):
        """ParseResult metadata should carry file-level information."""
        service = ParserService()
        result = await service.parse(sample_pdf_path, file_type="pdf")

        assert "file_type" in result.metadata
        assert result.metadata["file_type"] == "pdf"
        assert "filename" in result.metadata


class TestChunkAndParseResult:
    """Unit tests for the Chunk and ParseResult dataclasses."""

    def test_chunk_creation(self):
        """Chunk should store text, metadata, and index."""
        chunk = Chunk(
            text="几何学研究空间形状",
            metadata={"source": "test.pdf", "page": 1},
            chunk_index=0,
        )
        assert chunk.text == "几何学研究空间形状"
        assert chunk.metadata["source"] == "test.pdf"
        assert chunk.metadata["page"] == 1
        assert chunk.chunk_index == 0
