"""ParserService with support for PDF, Word, PPT, and image documents.

Handles Chinese-language educational materials using:
- PyMuPDF (fitz) for PDF
- python-docx for Word documents
- python-pptx for PowerPoint presentations
- Pillow for image metadata

Produces ParseResult with semantically chunked text and metadata per chunk.
"""

import logging
import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from app.core.exceptions import FileParsingError

logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class Chunk:
    """A single chunk of parsed text with metadata."""

    text: str
    metadata: dict = field(default_factory=dict)
    chunk_index: int = 0


@dataclass(frozen=True)
class ParseResult:
    """Result of parsing a document."""

    chunks: list[Chunk]
    page_count: int
    metadata: dict = field(default_factory=dict)


class ParserService:
    """Service for parsing educational materials into text chunks.

    Supports PDF, Word (docx), PowerPoint (pptx), and image files.
    Uses semantic chunking by splitting on section headers and paragraphs
    rather than arbitrary token counts.
    """

    SUPPORTED_TYPES: dict[str, str] = {
        "pdf": "pdf",
        "docx": "docx",
        "pptx": "pptx",
        "ppt": "pptx",
        "png": "image",
        "jpg": "image",
        "jpeg": "image",
        "gif": "image",
        "bmp": "image",
        "tiff": "image",
        "webp": "image",
    }

    async def parse(
        self,
        file_path: str,
        file_type: Optional[str] = None,
    ) -> ParseResult:
        """Parse a file and return structured chunks.

        Args:
            file_path: Path to the file to parse.
            file_type: Optional file type hint. If not provided,
                       detected from the file extension.

        Returns:
            ParseResult with chunks, page count, and file metadata.

        Raises:
            FileParsingError: If parsing fails.
            ValueError: If the file type is unsupported.
        """
        if not os.path.exists(file_path):
            raise FileParsingError(file_path, "File not found")

        detected_type = file_type or self._detect_type(file_path)
        if detected_type not in self.SUPPORTED_TYPES:
            raise ValueError(
                f"Unsupported file type: {detected_type}. "
                f"Supported: {', '.join(sorted(set(self.SUPPORTED_TYPES.values())))}"
            )

        filename = Path(file_path).name
        canonical_type = self.SUPPORTED_TYPES[detected_type]

        logger.info("Parsing %s as %s", filename, canonical_type)

        if canonical_type == "pdf":
            return await self._parse_pdf(file_path, filename)
        elif canonical_type == "docx":
            return await self._parse_docx(file_path, filename)
        elif canonical_type == "pptx":
            return await self._parse_pptx(file_path, filename)
        else:
            return await self._parse_image(file_path, filename)

    # ------------------------------------------------------------------
    # PDF parser
    # ------------------------------------------------------------------

    async def _parse_pdf(self, file_path: str, filename: str) -> ParseResult:
        """Parse PDF using PyMuPDF (fitz).

        Extracts text per page and creates chunks with page metadata.
        For semantic chunking, we split each page on paragraph boundaries
        and create one chunk per logical section.
        """
        try:
            import fitz

            doc = fitz.open(file_path)
            chunks: list[Chunk] = []
            chunk_idx = 0

            for page_num in range(doc.page_count):
                page = doc[page_num]
                page_text = page.get_text().strip()

                if not page_text:
                    continue

                # Split page text into paragraphs
                paragraphs = self._split_into_paragraphs(page_text)
                for para in paragraphs:
                    para = para.strip()
                    if not para:
                        continue
                    chunks.append(
                        Chunk(
                            text=para,
                            metadata={
                                "source": filename,
                                "page": page_num + 1,
                                "file_type": "pdf",
                            },
                            chunk_index=chunk_idx,
                        )
                    )
                    chunk_idx += 1

            page_count = doc.page_count
            doc.close()

            logger.info(
                "Parsed %s: %d pages, %d chunks", filename, page_count, len(chunks)
            )

            return ParseResult(
                chunks=chunks,
                page_count=page_count,
                metadata={
                    "filename": filename,
                    "file_type": "pdf",
                    "chunk_count": len(chunks),
                },
            )
        except Exception as exc:
            raise FileParsingError(filename, str(exc)) from exc

    # ------------------------------------------------------------------
    # Word parser
    # ------------------------------------------------------------------

    async def _parse_docx(self, file_path: str, filename: str) -> ParseResult:
        """Parse Word document using python-docx.

        Extracts paragraphs and creates chunks with section metadata.
        Headings are treated as section boundaries.
        """
        try:
            from docx import Document

            doc = Document(file_path)
            chunks: list[Chunk] = []
            chunk_idx = 0
            current_section: list[str] = []
            section_title = "Main"

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Heading detection: start a new chunk group
                if para.style.name.startswith("Heading"):
                    # Flush accumulated paragraphs as a chunk
                    if current_section:
                        combined = "\n".join(current_section)
                        chunks.append(
                            Chunk(
                                text=combined,
                                metadata={
                                    "source": filename,
                                    "section": section_title,
                                    "file_type": "docx",
                                },
                                chunk_index=chunk_idx,
                            )
                        )
                        chunk_idx += 1
                        current_section = []
                    section_title = text
                else:
                    current_section.append(text)

            # Flush remaining paragraphs
            if current_section:
                combined = "\n".join(current_section)
                chunks.append(
                    Chunk(
                        text=combined,
                        metadata={
                            "source": filename,
                            "section": section_title,
                            "file_type": "docx",
                        },
                        chunk_index=chunk_idx,
                    )
                )
                chunk_idx += 1

            # Estimate page count based on paragraph count
            page_count = max(1, len(doc.paragraphs) // 5)
            logger.info(
                "Parsed %s: %d pages (est.), %d chunks",
                filename,
                page_count,
                len(chunks),
            )

            return ParseResult(
                chunks=chunks,
                page_count=page_count,
                metadata={
                    "filename": filename,
                    "file_type": "docx",
                    "chunk_count": len(chunks),
                },
            )
        except Exception as exc:
            raise FileParsingError(filename, str(exc)) from exc

    # ------------------------------------------------------------------
    # PPT parser
    # ------------------------------------------------------------------

    async def _parse_pptx(self, file_path: str, filename: str) -> ParseResult:
        """Parse PowerPoint presentation using python-pptx.

        Extracts text per slide and creates chunks with slide metadata.
        """
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            chunks: list[Chunk] = []
            chunk_idx = 0
            slide_count = len(prs.slides)

            for slide_num, slide in enumerate(prs.slides, start=1):
                texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                texts.append(text)

                if not texts:
                    continue

                # Create one chunk per slide for semantic coherence
                combined = "\n".join(texts)
                chunks.append(
                    Chunk(
                        text=combined,
                        metadata={
                            "source": filename,
                            "slide": slide_num,
                            "file_type": "pptx",
                        },
                        chunk_index=chunk_idx,
                    )
                )
                chunk_idx += 1

            logger.info(
                "Parsed %s: %d slides, %d chunks",
                filename,
                slide_count,
                len(chunks),
            )

            return ParseResult(
                chunks=chunks,
                page_count=slide_count,
                metadata={
                    "filename": filename,
                    "file_type": "pptx",
                    "chunk_count": len(chunks),
                },
            )
        except Exception as exc:
            raise FileParsingError(filename, str(exc)) from exc

    # ------------------------------------------------------------------
    # Image parser
    # ------------------------------------------------------------------

    async def _parse_image(self, file_path: str, filename: str) -> ParseResult:
        """Parse image using Pillow.

        Extracts basic metadata (dimensions, format).
        Full OCR will be added in a future iteration.
        """
        try:
            from PIL import Image

            img = Image.open(file_path)
            width, height = img.size
            img_format = img.format or "unknown"

            # Create a single metadata chunk describing the image
            text_parts = [
                f"[Image: {filename}]",
                f"Format: {img_format}",
                f"Dimensions: {width}x{height}",
            ]
            text = "\n".join(text_parts)

            chunk = Chunk(
                text=text,
                metadata={
                    "source": filename,
                    "file_type": "image",
                    "width": width,
                    "height": height,
                    "format": img_format,
                },
                chunk_index=0,
            )

            logger.info("Parsed image %s: %dx%d %s", filename, width, height, img_format)

            return ParseResult(
                chunks=[chunk],
                page_count=1,
                metadata={
                    "filename": filename,
                    "file_type": "image",
                    "chunk_count": 1,
                },
            )
        except Exception as exc:
            raise FileParsingError(filename, str(exc)) from exc

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _detect_type(file_path: str) -> str:
        """Detect file type from extension."""
        ext = Path(file_path).suffix.lower().lstrip(".")
        supported_types = set(ParserService.SUPPORTED_TYPES.keys())
        if ext not in supported_types:
            raise ValueError(
                f"Unsupported file extension: .{ext}. "
                f"Supported: {', '.join(sorted(supported_types))}"
            )
        return ext

    @staticmethod
    def _split_into_paragraphs(text: str) -> list[str]:
        """Split text into paragraphs on blank lines or section markers.

        This produces semantically meaningful chunks rather than arbitrary
        token-level splits.
        """
        # Split on double newlines (blank lines)
        parts = text.split("\n\n")
        paragraphs: list[str] = []
        for part in parts:
            stripped = part.strip()
            if not stripped:
                continue
            # Further split long paragraphs at single newlines
            # but only if they are substantial
            if len(stripped) > 500:
                sub_parts = stripped.split("\n")
                current: list[str] = []
                current_len = 0
                for sub in sub_parts:
                    sub = sub.strip()
                    if not sub:
                        continue
                    current.append(sub)
                    current_len += len(sub)
                    if current_len > 300:
                        paragraphs.append("\n".join(current))
                        current = []
                        current_len = 0
                if current:
                    paragraphs.append("\n".join(current))
            else:
                paragraphs.append(stripped)

        return paragraphs